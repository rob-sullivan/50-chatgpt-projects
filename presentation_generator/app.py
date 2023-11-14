import collections.abc
import config
assert collections
import tkinter as tk
from pptx import Presentation
from pptx.util import Inches, Pt
import openai
from io import BytesIO
import requests
# API Token
openai.api_key = config.API_KEY

def slide_generator(text, prs):
    # I take text from the user and formulate it into a prompt for chatgpt
    prompt = f"Summarize the following text to a DALL-E image generation prompt: \n {text}"

    # I define which gpt I want to use then I pass it into chatgpt
    # Here I want the result to be summerized as a prompt for dalle
    model_engine = "gpt-4"
    dlp = openai.ChatCompletion.create(
        model=model_engine,
        messages=[
            {"role": "user", "content": "I will ask you a question"},
            {"role": "assistant", "content": "Ok"},
            {"role": "user", "content": f"{prompt}"}
        ],
        max_tokens=250,
        n=1,
        stop=None,
        temperature=0.8
    )
    # I define a dalle prompt and get output of chatgpt to pass into dalle
    dalle_prompt = dlp["choices"][0]["message"]["content"]

    # I ask dalle to generate an image based on chatgpt prompt
    response = openai.Image.create(
        prompt=dalle_prompt + " Style: digital art",
        n=1,#I only want one image
        size="1024x1024")
    image_url = response['data'][0]['url']

    #I ask chatGPT to create bullet points
    prompt = f"Create a bullet point text for a Powerpoint slide from the following text: \n {text}"
    ppt = openai.ChatCompletion.create(
        model=model_engine,
        messages=[
            {"role": "user", "content": "I will ask you a question"},
            {"role": "assistant", "content": "Ok"},
            {"role": "user", "content": f"{prompt}"}
        ],
        max_tokens=1024,
        n=1,
        stop=None,
        temperature=0.8
    )
    ppt_text = ppt["choices"][0]["message"]["content"]

    # I then ask chatgpt to add a title
    prompt = f"Create a title for a Powerpoint slide from the following text: \n {text}"
    ppt = openai.ChatCompletion.create(
        model=model_engine,
        messages=[
            {"role": "user", "content": "I will ask you a question"},
            {"role": "assistant", "content": "Ok"},
            {"role": "user", "content": f"{prompt}"}
        ],
        max_tokens=1024,
        n=1,
        stop=None,
        temperature=0.8
    )
    ppt_header = ppt["choices"][0]["message"]["content"]
    
    # I create a new slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    #Got image from url, convert it to bytes then add it to slide
    response = requests.get(image_url)
    img_bytes = BytesIO(response.content)
    slide.shapes.add_picture(img_bytes, Inches(1), Inches(1))

    # Added text box
    txBox = slide.shapes.add_textbox(Inches(3), Inches(1), Inches(4), Inches(1.5))
    tf = txBox.text_frame
    tf.text = ppt_text
    
    # added title
    title_shape = slide.shapes.title
    title_shape.text = ppt_header

def get_slides():
    # I get content from the text field starting from the first character to the last character, except the new line character.
    text = text_field.get("1.0", "end-1c")

    # I split text into paragraphs
    paragraphs = text.split("\n\n")

    # I initalise an empty powerpoint presentation
    prs = Presentation()
    width = Pt(1920)
    height = Pt(1080)
    prs.slide_width = width
    prs.slide_height = height

    # I loop through each text field paragraph and add them to the slides
    for paragraph in paragraphs:
        slide_generator(paragraph, prs)
        print("created slide")
    # Save with file name
    prs.save("chatgpt_presentation.pptx")

app = tk.Tk()
app.title("PPT Slides Generator")
app.geometry("800x600")

# Create text field
text_field = tk.Text(app)
text_field.pack(fill="both", expand=True)

text_field.configure(wrap="word", font=("Arial", 12))
text_field.focus_set()

# Create the button to create slides
create_button = tk.Button(app, text="Create Slides", command=get_slides)
create_button.pack()

app.mainloop()